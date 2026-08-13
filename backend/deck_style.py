"""
deck_style.py — modern dark-theme slide builder for the Lecture Assistant.

Drop this next to backend/main.py and replace the body of /download-pptx
with a call to build_deck(slides_data) -> file_path (see bottom of file
for the exact endpoint swap).

Design language (matches the reference decks the palette was picked for):
- Deep navy/near-black backgrounds throughout (no white slides).
- Small uppercase "kicker" label + page counter in the top corners of
  every slide — the one repeating structural motif, along with numbered
  icon badges for bullets (NOT an accent stripe/bar).
- One dominant color (navy), one sharp accent (amber), one supporting
  tone (cyan), one semantic warning tone (coral) reserved for Risks.
- Layout varies by slide role (title / intro / findings / risks /
  reading / generic) instead of repeating one template.
"""

import re
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn, nsmap

P_NS = 'http://schemas.openxmlformats.org/presentationml/2006/main'
P14_NS = 'http://schemas.microsoft.com/office/powerpoint/2010/main'

# ---------------------------------------------------------------------------
# Palette — "Modern Academic Indigo". Navy dominates (~65%), amber is the
# single sharp accent, cyan is the supporting tone, coral is reserved for
# Risks/warnings only so it keeps semantic meaning.
# ---------------------------------------------------------------------------
BG          = RGBColor(0x0A, 0x0E, 0x1B)   # near-black navy — slide background
CARD        = RGBColor(0x13, 0x19, 0x2C)   # card surface
CARD_ALT    = RGBColor(0x17, 0x1E, 0x35)   # slightly lighter card (badges, chips)
STROKE      = RGBColor(0x24, 0x2C, 0x47)   # hairline borders on cards
TEXT        = RGBColor(0xF3, 0xF1, 0xEA)   # primary text (off-white)
MUTED       = RGBColor(0x8B, 0x91, 0xA8)   # secondary / muted text
ACCENT      = RGBColor(0xF2, 0xB2, 0x0C)   # amber — primary accent
ACCENT_SOFT = RGBColor(0x3A, 0x30, 0x14)   # amber tint for badge fills
CYAN        = RGBColor(0x4D, 0xD6, 0xE0)   # supporting tone
CYAN_SOFT   = RGBColor(0x12, 0x2C, 0x30)
CORAL       = RGBColor(0xF2, 0x6B, 0x5E)   # warning tone — Risks only
CORAL_SOFT  = RGBColor(0x35, 0x18, 0x17)
LINK        = RGBColor(0x6E, 0xC9, 0xF2)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.6)

FONT_HEAD = "Calibri"     # safe-list sans, used for titles + body (renders true-to-width)
FONT_SERIF = "Cambria"    # safe-list serif accent, used sparingly for kickers/numerals


# ---------------------------------------------------------------------------
# Low-level helpers
# ---------------------------------------------------------------------------

def _no_autofit(tf):
    """Disable text autofit so our explicit font sizes are respected."""
    el = tf._txBody
    bodyPr = el.find(qn('a:bodyPr'))
    for tag in ('a:normAutofit', 'a:spAutoFit'):
        existing = bodyPr.find(qn(tag))
        if existing is not None:
            bodyPr.remove(existing)
    bodyPr.append(el.makeelement(qn('a:noAutofit'), {}))


def set_background(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill=None, line=None, line_w=Pt(0.75), shadow=False, radius=None):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius is not None else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, left, top, width, height)
    if radius is not None:
        try:
            shp.adjustments[0] = radius
        except Exception:
            pass
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = line_w
    shp.shadow.inherit = False
    if shadow:
        el = shp._element.spPr
        effectLst = el.makeelement(qn('a:effectLst'), {})
        shadow_el = el.makeelement(qn('a:outerShdw'), {
            'blurRad': '190500', 'dist': '38100', 'dir': '5400000', 'rotWithShape': '0'
        })
        clr = el.makeelement(qn('a:srgbClr'), {'val': '000000'})
        alpha = el.makeelement(qn('a:alpha'), {'val': '45000'})
        clr.append(alpha)
        shadow_el.append(clr)
        effectLst.append(shadow_el)
        el.append(effectLst)
    return shp


def add_circle(slide, left, top, diameter, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, diameter, diameter)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def add_text(slide, left, top, width, height, runs, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, line_spacing=1.0, space_after=0, wrap=True):
    """
    runs: list of paragraphs; each paragraph is a list of (text, size_pt, color, bold, font, italic) tuples,
          OR a plain string (rendered as one default run).
    Returns the created textbox.
    """
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    _no_autofit(tf)

    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        if isinstance(para, str):
            para = [(para, 14, TEXT, False, FONT_HEAD, False)]
        for spec in para:
            text, size, color, bold = spec[0], spec[1], spec[2], spec[3]
            font = spec[4] if len(spec) > 4 else FONT_HEAD
            italic = spec[5] if len(spec) > 5 else False
            link = spec[6] if len(spec) > 6 else None
            # Links no longer force an underline — a hyperlink run keeps
            # whatever color/weight the caller gives it (usually matched to
            # the surrounding text) so it blends in instead of looking like
            # a generic blue web link. Pass underline=True explicitly (9th
            # element) on the rare occasion one is actually wanted.
            underline = spec[7] if len(spec) > 7 else False
            r = p.add_run()
            r.text = text
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.italic = italic
            r.font.name = font
            r.font.color.rgb = color
            r.font.underline = underline
            if link:
                r.hyperlink.address = link
    return box


def linkify(text, size, color, bold=False, italic=False, font=FONT_HEAD):
    """Split text on inline URLs and return a single paragraph (list of run
    tuples) where the URL segments are made into real, clickable hyperlink
    runs. The link runs keep the exact same size/color/weight as the rest
    of the sentence — clickable, but visually blended into the body text
    rather than styled as a separate "blue underlined link"."""
    parts = [p for p in URL_RE.split(text) if p]
    runs = []
    for part in parts:
        if URL_RE.match(part):
            runs.append((part, size, color, bold, font, italic, part, False))
        else:
            runs.append((part, size, color, bold, font, italic, None, False))
    return runs or [(text, size, color, bold, font, italic, None, False)]


def citation_runs(citation, size=11, color=CYAN):
    """Render a trailing bullet citation (e.g. 'Source: Nature - https://...').
    If it contains a URL, the whole citation becomes one clickable hyperlink
    run — same citation color as always, no underline — so it reads as a
    normal citation that happens to be clickable, not a jarring web link."""
    urls = URL_RE.findall(citation)
    if urls:
        url = urls[0]
        label = URL_RE.sub('', citation).strip(' -\u2014:,') or url
        return [("\u2014 " + label, size, color, False, FONT_HEAD, True, url, False)]
    return [("\u2014 " + citation, size, color, False, FONT_HEAD, True)]


def add_slide_transition(slide, kind="fade", direction="l", speed="med"):
    """Attach a slide-to-slide transition (what PowerPoint calls
    'Transitions' — applies when advancing TO this slide).

    Fixed: the previous version wrapped fade/push in the p14 (PowerPoint
    2010+) extension namespace and added a p14:dur attribute. Neither is
    valid — per the OOXML/ISO-29500 schema, <p:fade/> and <p:push/> are
    plain base-schema elements (no extension needed), <p:transition> has
    no duration attribute at all (only spd="slow"/"med"/"fast"), and a
    schema violation here is exactly what makes PowerPoint's repair pass
    silently strip the transition — and the sibling <p:timing> animation
    block along with it, which is why neither was showing up."""
    sld = slide._element
    old = sld.find(qn('p:transition'))
    if old is not None:
        sld.remove(old)

    if kind == "push":
        inner = f'<p:push dir="{direction}"/>'
    else:
        inner = '<p:fade/>'

    xml = f'<p:transition xmlns:p="{P_NS}" spd="{speed}">{inner}</p:transition>'
    new_el = etree.fromstring(xml)

    cSld = sld.find(qn('p:cSld'))
    clrMapOvr = sld.find(qn('p:clrMapOvr'))
    anchor = clrMapOvr if clrMapOvr is not None else cSld
    anchor.addnext(new_el)


def add_click_build_animations(slide, shape_groups, duration_ms=400):
    """Attach a click-triggered 'build' animation to a slide: each group of
    shapes (e.g. a bullet's numbered badge + its text) fades in together on
    its own click, in the order the groups are given. This is the standard
    PowerPoint Animations-pane behavior (distinct from add_slide_transition,
    which only animates the change FROM the previous slide TO this one) —
    titles/cards/bullets build up progressively as the presenter advances
    instead of appearing all at once.

    Uses the OOXML p:timing tree directly since python-pptx has no native
    animation API. Written conservatively (Fade entrance, presetID 10) to
    match the structure PowerPoint itself generates for a basic 'Fade'
    build, since malformed timing XML can trigger a repair prompt on open.
    """
    shape_groups = [[s for s in g if s is not None] for g in shape_groups]
    shape_groups = [g for g in shape_groups if g]
    if not shape_groups:
        return

    sld = slide._element
    old = sld.find(qn('p:timing'))
    if old is not None:
        sld.remove(old)

    counter = {"n": 10}

    def nid():
        counter["n"] += 1
        return counter["n"]

    click_pars = []
    bld_entries = []
    for group in shape_groups:
        effect_pars = []
        for shp in group:
            sid = shp.shape_id
            preset_id = nid()
            beh_id = nid()
            effect_pars.append(
                f'<p:par><p:cTn id="{preset_id}" presetID="10" presetClass="entr" '
                f'presetSubtype="0" fill="hold" nodeType="clickEffect">'
                f'<p:stCondLst><p:cond delay="0"/></p:stCondLst>'
                f'<p:childTnLst><p:animEffect transition="in" filter="fade">'
                f'<p:cBhvr><p:cTn id="{beh_id}" dur="{duration_ms}"/>'
                f'<p:tgtEl><p:spTgt spid="{sid}"/></p:tgtEl></p:cBhvr>'
                f'</p:animEffect></p:childTnLst></p:cTn></p:par>'
            )
            bld_entries.append(f'<p:bldP spid="{sid}" grpId="0"/>')
        outer_id = nid()
        click_pars.append(
            f'<p:par><p:cTn id="{outer_id}" fill="hold">'
            f'<p:stCondLst><p:cond delay="indefinite"/></p:stCondLst>'
            f'<p:childTnLst>{"".join(effect_pars)}</p:childTnLst>'
            f'</p:cTn></p:par>'
        )

    main_id = nid()
    root_id = nid()
    xml = (
        f'<p:timing xmlns:p="{P_NS}"><p:tnLst><p:par>'
        f'<p:cTn id="{root_id}" dur="indefinite" restart="never" nodeType="tmRoot">'
        f'<p:childTnLst><p:seq concurrent="1" nextAc="seek">'
        f'<p:cTn id="{main_id}" dur="indefinite" nodeType="mainSeq">'
        f'<p:childTnLst>{"".join(click_pars)}</p:childTnLst></p:cTn>'
        f'<p:prevCondLst><p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>'
        f'<p:nextCondLst><p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>'
        f'</p:seq></p:childTnLst></p:cTn></p:par></p:tnLst>'
        f'<p:bldLst>{"".join(bld_entries)}</p:bldLst></p:timing>'
    )

    new_el = etree.fromstring(xml)
    sld.append(new_el)


def kicker(slide, label, index, total):
    """Top-left uppercase eyebrow label + top-right page counter. The one
    repeating navigational motif — present on every slide, never a stripe."""
    add_text(
        slide, MARGIN, Inches(0.42), Inches(7), Inches(0.35),
        [[(label.upper(), 11, ACCENT, True, FONT_HEAD)]],
        align=PP_ALIGN.LEFT,
    )
    add_text(
        slide, SLIDE_W - MARGIN - Inches(1.5), Inches(0.42), Inches(1.5), Inches(0.35),
        [[(f"{index + 1:02d} / {total:02d}", 11, MUTED, False, FONT_HEAD)]],
        align=PP_ALIGN.RIGHT,
    )
    # hairline under the header band
    ln = slide.shapes.add_connector(1, MARGIN, Inches(0.85), SLIDE_W - MARGIN, Inches(0.85))
    ln.line.color.rgb = STROKE
    ln.line.width = Pt(0.75)


def badge(slide, left, top, size, number_or_char, fill=ACCENT_SOFT, text_color=ACCENT, font_size=14):
    circle = add_circle(slide, left, top, size, fill)
    label = add_text(
        slide, left, top, size, size,
        [[(str(number_or_char), font_size, text_color, True, FONT_HEAD)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )
    return circle, label


CITATION_RE = re.compile(r'\s*[\[\(]([^\[\]\(\)]{2,60})[\]\)]\s*$')
URL_RE = re.compile(r'(https?://[^\s\]]+)')


def split_citation(text):
    """Pull a trailing [Author, Year] / (Source) citation off a bullet."""
    m = CITATION_RE.search(text.strip())
    if m:
        return text[:m.start()].strip(), m.group(1).strip()
    return text.strip(), None


# ---------------------------------------------------------------------------
# Slide role detection
# ---------------------------------------------------------------------------

def classify(title, index):
    t = title.lower()
    if index == 0 or 'title' in t and index == 0:
        return 'title'
    if 'introduc' in t:
        return 'intro'
    if 'risk' in t:
        return 'risks'
    if 'further reading' in t or 'reading' in t or 'references' in t:
        return 'reading'
    if 'key finding' in t or 'finding' in t:
        return 'findings'
    if 'summary' in t or 'conclusion' in t:
        return 'summary'
    return 'content'


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def build_title(slide, data, index, total):
    set_background(slide, BG)
    # soft glow accent tucked into the top-right corner — mostly ON the
    # slide, just a small bleed past the edge, not a shape hanging in space
    glow = add_circle(slide, SLIDE_W - Inches(2.6), Inches(-0.6), Inches(3.6), CARD_ALT)
    glow.fill.fore_color.rgb = CARD_ALT

    kicker_box = add_text(slide, MARGIN, Inches(0.6), Inches(6), Inches(0.35),
              [[("LECTURE BRIEF", 12, ACCENT, True)]])

    title = data.get('title', 'Untitled')
    title_box = add_text(slide, MARGIN, Inches(2.7), Inches(10.5), Inches(2.2),
              [[(title, 44, TEXT, True)]], line_spacing=1.05)

    groups = [[kicker_box], [title_box]]

    subtitle = data.get('subtitle', '')
    if subtitle:
        sub_box = add_text(slide, MARGIN, Inches(4.55), Inches(9), Inches(0.8),
                  [[(subtitle, 16, MUTED, False)]])
        groups.append([sub_box])

    add_text(slide, MARGIN, SLIDE_H - Inches(0.9), Inches(6), Inches(0.4),
              [[(f"{index + 1:02d} / {total:02d}", 11, MUTED, False)]])
    return groups


def build_intro(slide, data, index, total):
    set_background(slide, BG)
    kicker(slide, "Introduction", index, total)

    title_box = add_text(slide, MARGIN, Inches(1.25), Inches(10), Inches(0.5),
              [[(data.get('title', 'Introduction'), 28, TEXT, True)]])

    bullets = data.get('bullets', []) or ['']
    statement = bullets[0] if bullets else ''

    # oversized quote mark as the visual anchor
    quote_box = add_text(slide, MARGIN, Inches(2.5), Inches(1.5), Inches(1.2),
              [[("\u201C", 90, ACCENT_SOFT if False else RGBColor(0x2A, 0x30, 0x48), True, FONT_SERIF)]])

    statement_box = add_text(slide, Inches(1.4), Inches(3.0), Inches(10.6), Inches(2.2),
              [[(statement, 30, TEXT, True)]], line_spacing=1.2, anchor=MSO_ANCHOR.TOP)

    groups = [[title_box], [quote_box, statement_box]]

    subtitle = data.get('subtitle', '')
    if subtitle:
        sub_box = add_text(slide, Inches(1.4), Inches(5.3), Inches(9), Inches(0.6),
                  [[(subtitle, 14, MUTED, False, FONT_HEAD, True)]])
        groups.append([sub_box])
    return groups


def _bullet_row(slide, left, top, width, index_label, body_text, citation=None,
                 badge_fill=ACCENT_SOFT, badge_text=ACCENT, row_h=Inches(0.9)):
    b_circle, b_label = badge(slide, left, top, Inches(0.42), index_label, fill=badge_fill, text_color=badge_text, font_size=13)
    text_left = left + Inches(0.62)
    text_w = width - Inches(0.62)
    paras = [linkify(body_text, 15, TEXT)]
    if citation:
        paras.append(citation_runs(citation))
    text_box = add_text(slide, text_left, top - Inches(0.03), text_w, row_h, paras,
              line_spacing=1.15, space_after=2)
    return [b_circle, b_label, text_box]


def build_findings(slide, data, index, total):
    set_background(slide, BG)
    kicker(slide, "Key Findings", index, total)

    title_box = add_text(slide, MARGIN, Inches(1.15), Inches(10.5), Inches(0.6),
              [[(data.get('title', 'Key Findings'), 30, TEXT, True)]])
    groups = [[title_box]]
    subtitle = data.get('subtitle', '')
    if subtitle:
        sub_box = add_text(slide, MARGIN, Inches(1.75), Inches(10.5), Inches(0.4),
                  [[(subtitle, 13, MUTED, False)]])
        groups[-1].append(sub_box)

    bullets = data.get('bullets', [])
    card_top = Inches(2.35)
    card_h = SLIDE_H - card_top - Inches(0.55)
    card = add_rect(slide, MARGIN, card_top, SLIDE_W - 2 * MARGIN, card_h, fill=CARD, line=STROKE, radius=0.045)
    groups.append([card])

    pad = Inches(0.4)
    n = max(len(bullets), 1)
    row_h = (card_h - 2 * pad) / n
    y = card_top + pad - Inches(0.05)
    for i, b in enumerate(bullets):
        text, cite = split_citation(str(b))
        row_shapes = _bullet_row(slide, MARGIN + pad, y, SLIDE_W - 2 * MARGIN - 2 * pad,
                    str(i + 1), text, citation=cite, row_h=row_h)
        groups.append(row_shapes)
        y += row_h
    return groups


def build_risks(slide, data, index, total):
    set_background(slide, BG)
    kicker(slide, "Risks", index, total)

    title_box = add_text(slide, MARGIN, Inches(1.15), Inches(10.5), Inches(0.6),
              [[(data.get('title', 'Risks'), 30, TEXT, True)]])
    groups = [[title_box]]
    subtitle = data.get('subtitle', '')
    if subtitle:
        sub_box = add_text(slide, MARGIN, Inches(1.75), Inches(10.5), Inches(0.4),
                  [[(subtitle, 13, MUTED, False)]])
        groups[-1].append(sub_box)

    bullets = (data.get('bullets', []) or [])[:3]
    gap = Inches(0.35)
    n = max(len(bullets), 1)
    total_w = SLIDE_W - 2 * MARGIN - gap * (n - 1)
    col_w = total_w / n
    top = Inches(2.4)
    height = Inches(4.1)

    for i, b in enumerate(bullets):
        left = MARGIN + i * (col_w + gap)
        card = add_rect(slide, left, top, col_w, height, fill=CARD, line=STROKE, radius=0.06)
        b_circle, b_label = badge(slide, left + Inches(0.3), top + Inches(0.3), Inches(0.5),
              "!", fill=CORAL_SOFT, text_color=CORAL, font_size=18)
        text_box = add_text(slide, left + Inches(0.3), top + Inches(1.05), col_w - Inches(0.6), height - Inches(1.35),
                  [[(str(b), 15, TEXT, False)]], line_spacing=1.25)
        groups.append([card, b_circle, b_label, text_box])
    return groups


def build_reading(slide, data, index, total):
    set_background(slide, BG)
    kicker(slide, "Further Reading", index, total)

    title_box = add_text(slide, MARGIN, Inches(1.15), Inches(10.5), Inches(0.6),
              [[(data.get('title', 'Further Reading'), 30, TEXT, True)]])
    groups = [[title_box]]

    items = data.get('bullets', [])
    left_col = items[:len(items) // 2 + len(items) % 2]
    right_col = items[len(left_col):]

    col_w = (SLIDE_W - 2 * MARGIN - Inches(0.5)) / 2
    top0 = Inches(2.05)
    avail_h = SLIDE_H - top0 - Inches(0.5)

    def render_col(col_items, left, start_num):
        n = max(len(col_items), 1)
        row_h = avail_h / n
        y = top0
        for i, item in enumerate(col_items):
            item_str = str(item)
            urls = URL_RE.findall(item_str)
            display = URL_RE.sub('', item_str).strip(' -\u2014')
            if not display:
                display = urls[0] if urls else item_str

            b_circle, b_label = badge(slide, left, y + Inches(0.05), Inches(0.4), start_num + i,
                  fill=CYAN_SOFT, text_color=CYAN, font_size=12)
            # The whole title/reference text becomes the clickable link when
            # a URL is present — same off-white color as any other item, no
            # separate URL line, no underline, so it reads as plain text
            # that happens to be clickable rather than a distinct "link".
            link = urls[0] if urls else None
            paras = [[(display, 14, TEXT, False, FONT_HEAD, False, link, False)]]
            text_box = add_text(slide, left + Inches(0.58), y, col_w - Inches(0.58), row_h,
                      paras, line_spacing=1.1, space_after=2)
            groups.append([b_circle, b_label, text_box])
            y += row_h

    render_col(left_col, MARGIN, 1)
    render_col(right_col, MARGIN + col_w + Inches(0.5), len(left_col) + 1)
    return groups


def build_content(slide, data, index, total, role_label="Overview"):
    """Generic layout: left title/subtitle column, right bullet card.
    Used for Summary and any other slide the plan introduces."""
    set_background(slide, BG)
    kicker(slide, role_label, index, total)

    left_w = Inches(4.5)
    title_box = add_text(slide, MARGIN, Inches(2.0), left_w, Inches(2.2),
              [[(data.get('title', 'Untitled'), 30, TEXT, True)]], line_spacing=1.05)
    groups = [[title_box]]
    subtitle = data.get('subtitle', '')
    if subtitle:
        sub_box = add_text(slide, MARGIN, Inches(3.5) if len(data.get('title', '')) < 30 else Inches(4.0),
                  left_w - Inches(0.2), Inches(1.5),
                  [[(subtitle, 14, ACCENT, False, FONT_HEAD, True)]], line_spacing=1.2)
        groups.append([sub_box])

    bullets = data.get('bullets', [])
    card_left = MARGIN + left_w + Inches(0.4)
    card_top = Inches(1.15)
    card_w = SLIDE_W - card_left - MARGIN
    card_h = SLIDE_H - card_top - Inches(0.55)
    card = add_rect(slide, card_left, card_top, card_w, card_h, fill=CARD, line=STROKE, radius=0.04)
    groups.append([card])

    pad = Inches(0.4)
    n = max(len(bullets), 1)
    row_h = (card_h - 2 * pad) / n
    y = card_top + pad - Inches(0.05)
    for i, b in enumerate(bullets):
        text, cite = split_citation(str(b))
        row_shapes = _bullet_row(slide, card_left + pad, y, card_w - 2 * pad,
                    str(i + 1), text, citation=cite, row_h=row_h)
        groups.append(row_shapes)
        y += row_h
    return groups


ROLE_LABELS = {
    'title': 'Lecture Brief',
    'intro': 'Introduction',
    'findings': 'Key Findings',
    'risks': 'Risks',
    'reading': 'Further Reading',
    'summary': 'Summary',
    'content': 'Overview',
}


# ---------------------------------------------------------------------------
# Entry point
# ---------------------------------------------------------------------------

def build_deck(slides_data, out_path="lecture_slides_temp.pptx"):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank_layout = prs.slide_layouts[6]

    total = len(slides_data)
    for index, data in enumerate(slides_data):
        slide = prs.slides.add_slide(blank_layout)
        role = classify(data.get('title', ''), index)
        label = ROLE_LABELS.get(role, 'Overview')

        if role == 'title':
            build_title(slide, data, index, total)
        elif role == 'intro':
            build_intro(slide, data, index, total)
        elif role == 'findings':
            build_findings(slide, data, index, total)
        elif role == 'risks':
            build_risks(slide, data, index, total)
        elif role == 'reading':
            build_reading(slide, data, index, total)
        else:
            build_content(slide, data, index, total, role_label=label)

        # Slide-to-slide transition (fade). Applied to every slide so
        # advancing the deck never feels like a hard cut.
        add_slide_transition(slide, kind="fade", speed="med")

        # NOTE: per-element "click to build" animations (title/cards/bullets
        # fading in one at a time) were removed — the hand-written OOXML
        # timing XML validated correctly against the schema, but in
        # practice advancing to the next click made shapes blink instead
        # of fading. Getting PowerPoint's animation timing exactly right
        # without a real PowerPoint instance to test against isn't
        # reliable, so it's off rather than shipping something flaky.
        # The builder functions still return their shapes in build order
        # (`groups`, currently unused here) if this is worth revisiting.

    prs.save(out_path)
    return out_path


if __name__ == "__main__":
    # Quick local smoke test with dummy data shaped like the LLM's JSON output.
    sample = [
        {"title": "Introduction to Quantum Error Correction", "subtitle": "Context and background",
         "bullets": []},
        {"title": "Introduction", "subtitle": "Why this matters",
         "bullets": ["Quantum error correction protects fragile qubit information from decoherence by encoding one logical qubit across many physical qubits."]},
        {"title": "Core Concepts", "subtitle": "The fundamentals",
         "bullets": [
             "The threshold theorem shows that if physical error rates fall below roughly 1%, arbitrarily long quantum computation becomes possible through recursive encoding [Aharonov & Ben-Or, 1997].",
             "Surface codes arrange physical qubits on a 2D lattice and detect errors via repeated stabilizer measurements without collapsing the encoded state [Fowler et al., 2012].",
             "Logical error rates drop exponentially as code distance increases, at the cost of a roughly quadratic increase in physical qubit count.",
             "Real hardware from Google and IBM has demonstrated below-threshold error suppression as code distance scales from 3 to 5 [Google Quantum AI, 2023].",
         ]},
        {"title": "Key Findings", "subtitle": "What the data shows",
         "bullets": [
             "Surface-code experiments crossed the break-even point in 2023, with logical qubits outliving their best physical qubit for the first time [Google Quantum AI, 2023].",
             "Error rates below 0.1% per gate are now achievable on superconducting hardware using dynamical decoupling and improved calibration [IBM, 2024].",
             "Cross-platform comparisons show trapped-ion systems reaching higher gate fidelities but lower qubit counts than superconducting systems [Quantinuum, 2023].",
             "Decoder latency remains a bottleneck for real-time correction, with most decoders still too slow for fault-tolerant clock speeds [Delfosse, 2020].",
         ]},
        {"title": "Risks", "subtitle": "What could go wrong",
         "bullets": [
             "Correlated errors from cosmic rays can strike many physical qubits simultaneously, defeating the independence assumption most codes rely on.",
             "Decoder complexity grows rapidly with code distance, risking a mismatch between classical control speed and qubit coherence time.",
             "Overreliance on near-term error mitigation could delay investment in true fault tolerance, stalling long-term scaling roadmaps.",
         ]},
        {"title": "Summary", "subtitle": "Where things stand",
         "bullets": [
             "Break-even error correction is now experimentally real, not just theoretical, marking a turning point for the field.",
             "Surface codes remain the leading near-term architecture due to their high error threshold and local connectivity requirements.",
             "Decoder speed and correlated-error resilience are now the binding constraints on the path to fault tolerance.",
         ]},
        {"title": "Further Reading", "subtitle": "",
         "bullets": [
             "Fowler, Mariantoni, Martinis & Cleland, 'Surface codes: Towards practical large-scale quantum computation' https://arxiv.org/abs/1208.0928",
             "Google Quantum AI, 'Suppressing quantum errors by scaling a surface code logical qubit' (Nature, 2023)",
             "Preskill, 'Quantum Computing in the NISQ era and beyond' https://arxiv.org/abs/1801.00862",
             "Nielsen & Chuang, Quantum Computation and Quantum Information (Cambridge University Press)",
             "Terhal, 'Quantum error correction for quantum memories' https://arxiv.org/abs/1302.3428",
             "IBM Quantum, 'Error correction roadmap' ibm.com/quantum/roadmap",
         ]},
    ]
    build_deck(sample, "/home/claude/build/demo_deck.pptx")
    print("done")